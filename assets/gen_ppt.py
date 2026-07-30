# -*- coding: utf-8 -*-
"""
FX Sentinel — 기술설명서 PPTX 생성기
발표자료.md(17장 원고)를 16:9 슬라이드로. 재생성: python assets/gen_ppt.py
출력: FX_Sentinel_기술설명서.pptx (temp/KB 루트)
"""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
A = lambda n: os.path.join(HERE, n)

NAVY  = RGBColor(0x1A, 0x2B, 0x4A)
GOLD  = RGBColor(0xFF, 0xBC, 0x00)
INK   = RGBColor(0x22, 0x31, 0x4E)
GRAY  = RGBColor(0x6B, 0x76, 0x86)
BG    = RGBColor(0xF7, 0xF8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED   = RGBColor(0xC0, 0x39, 0x2B)
GOOD  = RGBColor(0x1E, 0x7A, 0x46)
FONT  = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
PAGE_W, PAGE_H = 13.333, 7.5
MARG = 0.55

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh

def runs_md(p, text, size, color, bold=False, bold_color=None):
    """**볼드** 마크업 지원 런 생성."""
    for i, seg in enumerate(re.split(r"\*\*", text)):
        if seg == "": continue
        r = p.add_run(); r.text = seg
        r.font.name = FONT; r.font.size = Pt(size)
        b = bold or (i % 2 == 1)
        r.font.bold = b
        r.font.color.rgb = (bold_color or color) if (i % 2 == 1 and bold_color) else color
        if b and not bold and bold_color is None:
            r.font.color.rgb = NAVY if color in (INK, GRAY) else color

def txt(s, x, y, w, h, lines, size=15, color=INK, bold=False, align=PP_ALIGN.LEFT,
        leading=1.18, space_after=7, bold_color=None, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = leading; p.space_after = Pt(space_after)
        runs_md(p, ln, size, color, bold=bold, bold_color=bold_color)
    return tb

def bullets(s, x, y, w, h, items, size=14.5, gap=8, color=INK):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.2; p.space_after = Pt(gap)
        r = p.add_run(); r.text = "▎ "; r.font.name = FONT; r.font.size = Pt(size)
        r.font.bold = True; r.font.color.rgb = GOLD
        runs_md(p, it, size, color)
    return tb

def headline(s, kicker, title, sub=None):
    rect(s, 0, 0, PAGE_W, 0.06, GOLD)
    txt(s, MARG, 0.34, 12.2, 0.32, kicker, size=12.5, color=GRAY, bold=True)
    txt(s, MARG, 0.66, 12.4, 0.7, title, size=25, color=NAVY, bold=True, space_after=0)
    if sub:
        txt(s, MARG, 1.28, 12.3, 0.4, sub, size=13.5, color=GRAY, space_after=0)

_pageno = [0]
def footer(s, note=None):
    _pageno[0] += 1
    txt(s, MARG, 7.08, 9.5, 0.3, note or "FX Sentinel — 제8회 KB AI Challenge · 현직자 Pick #6 수출입 금융 지원 에이전트",
        size=9.5, color=GRAY, space_after=0)
    txt(s, 12.35, 7.08, 0.5, 0.3, str(_pageno[0]), size=10, color=GRAY, align=PP_ALIGN.RIGHT, space_after=0)

def picture(s, path, x, y, w, border=True):
    from PIL import Image
    iw, ih = Image.open(path).size
    h = w * ih / iw
    if border:
        rect(s, x - 0.03, y - 0.03, w + 0.06, h + 0.06, RGBColor(0xE2, 0xE6, 0xEC))
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w))
    return h

def table(s, x, y, w, rows, col_w, header=True, size=12.5, row_h=0.34):
    shp = s.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y),
                             Inches(w), Inches(row_h * len(rows)))
    t = shp.table
    for ci, cw in enumerate(col_w):
        t.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, cell in enumerate(row):
            c = t.cell(ri, ci)
            c.margin_left = Inches(0.08); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            if header and ri == 0:
                c.fill.fore_color.rgb = NAVY
            else:
                c.fill.fore_color.rgb = WHITE if ri % 2 == 1 else BG
            tfr = c.text_frame; tfr.word_wrap = True
            p = tfr.paragraphs[0]; p.line_spacing = 1.08
            col = WHITE if (header and ri == 0) else INK
            runs_md(p, str(cell), size, col, bold=(header and ri == 0))
    return shp

# ══════════════════ 1. 표지 ══════════════════
s = slide()
rect(s, 0, 0, PAGE_W, PAGE_H, NAVY)
rect(s, 0, 4.92, PAGE_W, 0.045, GOLD)
txt(s, 1.0, 1.05, 11.3, 0.4, "제8회 KB AI Challenge (2026) · 현직자 Pick #6 「수출입 금융 지원 에이전트」",
    size=14, color=GOLD, bold=True)
txt(s, 1.0, 1.75, 11.3, 1.2, "FX Sentinel", size=58, color=WHITE, bold=True)
txt(s, 1.0, 2.95, 11.3, 0.6, "KB Star FX 기반 수출입 금융 AI 코파일럿", size=24, color=WHITE)
txt(s, 1.0, 3.75, 11.3, 0.8,
    ["환율을 맞히는 서비스가 아닙니다.",
     "수출입 기업의 **예산환율이 깨질 확률(BBP)**을 계산해, KB국민은행의 실제 외환·수출입금융·무역지원 상품으로 안전하게 라우팅합니다."],
    size=15.5, color=RGBColor(0xC9, 0xD2, 0xE2), bold_color=GOLD, space_after=5)
txt(s, 1.0, 5.25, 11.3, 0.4, "▶ 라이브 데모 · https://seonu-dragon.github.io/KB_FX_Sentinel/",
    size=16, color=GOLD, bold=True)
txt(s, 1.0, 5.75, 11.3, 0.4, "오프라인 단일 파일(FX_Sentinel_demo_ui.html)로도 동일 화면 · 외부 의존성 0",
    size=12.5, color=RGBColor(0x9A, 0xA7, 0xBD))

# ══════════════════ 2. 문제 정의 ══════════════════
s = slide()
headline(s, "문제 정의 — 현직자 Pick #6 · 수출입 금융 지원", "“예산환율은 정해뒀지만, 그게 깨질 확률은 아무도 모른다.”")
bullets(s, MARG, 1.75, 12.2, 4.0, [
 "중소 수출입기업은 손익분기 환율(예산환율)을 정해두지만, **깨질 확률을 정량으로 아는 곳은 드물다.**",
 "환헤지는 “할까/말까” 이분법 — **얼마를 · 어떤 수단으로** 헤지할지 근거가 없다.",
 "결제방식(T/T·L/C·D/A·O/A)·거래국·품목마다 필요한 KB 상품이 다른데 **어디에 뭘 신청할지 정보 비대칭**이 크다.",
 "개인사업자·글로벌셀러는 정산 환율·수수료·환전 타이밍에 특히 취약.",
], size=17, gap=16)
rect(s, MARG, 5.5, 12.2, 1.1, BG)
txt(s, 0.85, 5.68, 11.6, 0.8,
    "타깃 — 50~60대 SME 재무담당자. “전문 트레이더 도구”가 아니라 **사장님이 한눈에 이해하는 리스크 번역기.**",
    size=15.5, color=INK)
footer(s)

# ══════════════════ 3. 왜 KB인가 ══════════════════
s = slide()
headline(s, "왜 KB인가", "맨땅 신규가 아니라, KB Star FX 위에 얹는 증강 레이어")
table(s, MARG, 1.8, 12.2, [
 ["KB가 이미 가진 것", "비어 있는 칸 (= FX Sentinel)"],
 ["AI 환율전망 (뉴스감성 컨센서스)", "회사별 **예산환율 이탈확률 BBP** — 정량·개인화·확률"],
 ["헤지솔루션 · 환관리센터 · 경제캘린더", "결제방식·거래국·품목 → **KB 상품 라우팅 오케스트레이션**"],
], [6.1, 6.1], size=14.5, row_h=0.62)
bullets(s, MARG, 4.1, 12.2, 2.0, [
 "KB가 채운 축은 **뉴스심리(S) 한 축** — 그 위에 **정량 조기경보(BBP) + 개인화 + KB 공식상품 라우팅**을 얹는다.",
 "**다른 은행 로고로 대체 불가** — KB Payment Usance · 특별출연(K-SURE) · 해외거래처 신용조사 · 파트너관세사가 라우팅의 목적지.",
 "배포 채널(KB Star FX)이 이미 존재 → **실배포 논리가 곧바로 선다.**",
], size=15.5, gap=13)
footer(s)

# ══════════════════ 4. 솔루션 개요 ══════════════════
s = slide()
headline(s, "솔루션 개요 — 데모와 1:1", "3스텝 사용자 여정 — 화면 그대로")
txt(s, MARG, 1.62, 12.2, 0.5,
    "**① 거래 정보 입력**(간편 5개 + 인보이스 업로드 = 실수요 게이트)  →  **② 리스크 진단**(BBP·의사결정 도우미·XAI)  →  **③ 금융지원 라우팅**(KB 상품 → RM 티켓)",
    size=14, color=INK)
picture(s, A("ppt_core_3step.jpg"), 2.37, 2.2, 8.6)
footer(s)

# ══════════════════ 5. 데모① ══════════════════
s = slide()
headline(s, "[데모 ①] 라이브 — 나래상사(수입) End-to-End ★", "한 화면에서 경보 → 진단 → 근거 → 상품 → 티켓")
picture(s, A("ppt_home_hero.jpg"), MARG, 1.7, 8.3)
bullets(s, 9.15, 1.8, 3.7, 4.9, [
 "열자마자 **에이전트 점검 브리핑** — “이미 점검해 두었습니다”를 1인칭 보고",
 "판정 배너 **초과 가능성 64.3%** · 부담액(보수적 상한) 1,833만원",
 "**4-에이전트 파이프라인**이 실계산값으로 순차 점등",
 "**자가 반박** — 왜 100%가 아니라 부분 50%인지 스스로 설명",
 "`상담 예약 요청` → **구조화 RM 티켓**(상품·서류·확인사항·감사근거)",
], size=12.5, gap=9)
footer(s)

# ══════════════════ 6. 데모② 개인화 ══════════════════
s = slide()
headline(s, "[데모 ②] 개인화 실증", "같은 시장, 다른 회사 → 다른 BBP · 다른 KB 상품")
table(s, MARG, 1.75, 12.2, [
 ["프로필", "거래", "BBP", "1순위 KB 상품"],
 ["한빛정밀 · 수출확정", "수출 · 미국 · L/C", "16.5%", "수출환어음매입 (L/C Nego)"],
 ["대성무역 · 가결제", "수출 · 베트남 · PO단계", "26.9%", "무역금융(수출) · 선물환 차단(실수요 미확정)"],
 ["소망전자 · 무여신", "수출 · 중국 · O/A", "32.7%", "KB 특별출연 수출입 금융지원 (K-SURE)"],
 ["나래상사 · 수입 (기본)", "수입 · 베트남 · T/T", "64.3%", "KB Payment Usance"],
 ["하늘샵 · 글로벌셀러", "개인사업자 · 미국 · O/A", "1.2%", "KB 글로벌셀러 우대서비스"],
], [2.6, 2.9, 1.3, 5.4], size=13, row_h=0.5)
bullets(s, MARG, 5.15, 12.2, 1.4, [
 "프리셋을 눌러가며 BBP·1순위 상품이 실시간으로 갈린다 = **개인화 리스크 번역**의 실증 (하드코딩 아님 — 폼 입력에서 파생)",
 "시장 국면을 **2022 달러 강세**로 바꾸면 전 지표·큐가 재계산 — **급변기 대응**이 핵심 가치",
], size=14.5, gap=10)
footer(s)

# ══════════════════ 7. 데모③ 소비자보호 ══════════════════
s = slide()
headline(s, "[데모 ③] 소비자보호 — 금소법 판매 게이트", "AI가 상품을 더 파는 게 아니라, 팔면 안 되는 상대를 스스로 거른다")
picture(s, A("ppt_suitability.jpg"), MARG, 1.7, 8.3)
bullets(s, 9.15, 1.8, 3.7, 4.9, [
 "소망전자 → **적정성 2/5 부적정**(손익구조 이해·감내 손실 미충족)",
 "범위선물환 **[권유보류]** — AI가 자진 철회",
 "단 **K-SURE 환변동보험·특별출연은 남긴다** — 규제로 포용을 죽이지 않는다",
 "여신+파생 동시추천 → **꺾기(구속성) 점검 플래그**",
 "금소법 6대 판매원칙 매핑 · 문항·임계는 데모 설계값임을 화면에 명시",
], size=12.5, gap=9)
footer(s)

# ══════════════════ 8. 데모④ KIKO ══════════════════
s = slide()
headline(s, "[데모 ④] 은행 리스크 — KIKO를 계약 전에 재현", "계약 후 터질 손실을, 계약 전에 시뮬레이션한다")
picture(s, A("ppt_kiko_stress.jpg"), MARG, 1.7, 8.3)
bullets(s, 9.15, 1.8, 3.7, 4.9, [
 "체결 300k 선물환에 **2022 달러 강세 국면** 적용",
 "방향을 고정하지 않는다 — 수입·매수 포지션의 아픈 방향은 **원화 강세**",
 "−2σ: **평가손 7,450만 · 추가담보 2,863만 → 부족 1,363만 (감당 불가)**",
 "**“체결 명목 153,307 이하였어야”**를 역산 제시",
 "KIKO 경로(평가손→추가담보→자금경색)를 계약 전에 차단",
], size=12.5, gap=9)
footer(s)

# ══════════════════ 9. BBP 수식 ══════════════════
s = slide()
headline(s, "자체 정량 지표 — BBP (요건 #1)", "“GARCH 재포장”이 아니라, 회사별 손익분기가 깨질 확률")
rect(s, MARG, 1.75, 12.2, 1.15, NAVY)
txt(s, 0.9, 1.98, 11.6, 0.7,
    "BBP_t  =  P( S(t+N) 가 예산환율 K 를 불리하게 이탈 )  =  T_ν( (ln K − ln S_t) / (σ_fwd · √(N/252)) )",
    size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, MARG, 3.25, 12.2, 3.3, [
 "**개인화** — 회사의 K·포지션·만기 N → “변동성 숫자”가 아니라 **이 회사가 깨질 확률**",
 "**팻테일** — Student-t(ν=5)로 외환 꼬리 반영 (정규분포의 과소평가 교정)",
 "**정책개입 보정** — 라운드넘버(1,300·1,400) + 개입 국면 비대칭 절단",
 "**무예측 정직성** — drift=0(랜덤워크): 방향을 예측하지 않고, 이미 정해진 포지션에서 결정론적으로 도출",
 "코어는 `fx_sentinel/bbp.py`(Python 실측) — 데모 JS가 동일 수식 포팅 · **자체 IP**",
], size=15, gap=13)
footer(s)

# ══════════════════ 10. 캘리브레이션 ══════════════════
s = slide()
headline(s, "BBP 캘리브레이션 — 정직 고지 (숫자 옆에서 말한다)", "신뢰성은 화면에서 증명하고, 한계도 화면에서 고지한다")
picture(s, A("calibration_chart.png"), MARG, 1.75, 6.1)
bullets(s, 7.0, 1.9, 5.8, 4.6, [
 "만기별 ECE — 1M 0.036 / **3M 0.034** / 6M 0.084 · “BBP 65% → 실제 66%” 수준 정합",
 "2018–2026 실데이터 **3,800건 무누수 백테스트** (Brier 0.175)",
 "**한계 동시 고지** — 방향별로 나누면 ECE 0.139(4배) · 수입 실제 초과빈도 74.5% · ES는 실측의 1.86배 보수적",
 "→ 화면 표기는 “예상 부담액”이 아니라 **“부담액(보수적 상한)”** + 방향별 오차를 함께 고지",
 "출처: `state/bbp_calibration.json` · `bbp_validation.json` — 화면 수치와 1:1",
], size=13.5, gap=11)
footer(s)

# ══════════════════ 11. 경제성 ══════════════════
s = slide()
headline(s, "경제성 — 부분헤지 프론티어", "체계적 부분헤지가 무헤지·상시헤지 양극단을 이긴다 (리스크·비용·꼬리 3면)")
picture(s, A("frontier_kr.png"), 2.17, 1.75, 9.0)
bullets(s, MARG, 5.55, 12.2, 1.3, [
 "수출 SME(2019–2026) 하방분산 — 무헤지 759 · **부분 50%(권장) 460** · 상시 100% 571 · 비용은 상시의 절반",
 "**삭제한 과대주장** — “AI가 동적 타이밍을 잘한다”는 walk-forward 반증으로 폐기 · “AI 마법”이 아니라 **규율 있는 부분헤지**",
], size=14, gap=9)
footer(s)

# ══════════════════ 12. 정직 검증 ══════════════════
s = slide()
headline(s, "정직 검증 = 차별화 (요건 #6)", "대부분 팀이 안 하는 지적 정직성을 전면에")
picture(s, A("ppt_governance.jpg"), MARG, 1.7, 8.3)
bullets(s, 9.15, 1.8, 3.7, 4.9, [
 "**walk-forward OOS · PBO 0.20 · DSR 0.9999**를 실제로 돌려 예측력 측정",
 "Phase1 AUC 0.577 < 레벨 0.596 · 동적헤지 OOS 패배 → **예측 알파를 스스로 기각**",
 "낡은 주장을 몰래 지우지 않는다 — “검증했고, 안 통하는 건 걷어냈다”",
 "근거·거버넌스·로드맵은 데모 **④ 검증·거버넌스 탭**에 집결 (RM 실무 화면과 분리)",
], size=12.5, gap=10)
footer(s)

# ══════════════════ 13. 에이전트란 ══════════════════
s = slide()
headline(s, "에이전트란 무엇인가 (요건 #3 정면 대응) ★", "챗봇은 인터페이스일 뿐 — 에이전트는 “묻기 전에 일하는 주체”다")
table(s, MARG, 1.7, 5.9, [
 ["루프", "FX Sentinel의 실물"],
 ["인지", "시세·국면(σ·EWI)·프로필·서류(OCR)·자연어"],
 ["추론", "BBP 엔진·XAI 분해·룰엔진(자격·적정성·한도)"],
 ["행동", "선제 경보·상품 라우팅·RM 티켓·감시 알림"],
 ["감독", "준법 게이트·RM 최종 확인·영업점 계약 확정"],
], [1.1, 4.8], size=12, row_h=0.44)
bullets(s, MARG, 4.2, 5.9, 2.3, [
 "**AI 모드(채팅)** — 현재 진단 근거로 종합 답변·기업 전환",
 "예측 질문(“내일 환율 오를까?”)은 **“범위 밖 — 답을 지어내지 않습니다”** 거절 → RM 안내",
 "문장 입력·감시 위임·Advisor 질문 — 전부 행내 처리 · 외부 전송 0",
], size=12.5, gap=8)
picture(s, A("ppt_ai_mode.jpg"), 6.85, 1.7, 6.0)
footer(s)

# ══════════════════ 14. 멀티에이전트 + 정보보호 ══════════════════
s = slide()
headline(s, "멀티에이전트 구조 + 정보보호 (요건 #3·#4)", "경보 → 설명 → 정보 → 번역 → 제안 → 티켓 · 고객정보는 밖으로 안 나간다")
picture(s, A("agent_flow.png"), MARG, 1.7, 12.2)
bullets(s, MARG, 5.15, 12.2, 1.8, [
 "**Sentinel**(결정론 코어·경보) → **Analyst**(V/J/C XAI·LLM 설명층+규칙 폴백) → **Advisor**(지식베이스·출처·가드레일) → **Hedge**(CIP·트리아지·KB상품)",
 "고객 화면 AI 요약은 **행내(브라우저) 생성** — 숫자는 결정론 엔진만 만들므로 **환각 원천 차단** · 기준환율(영업비밀)·신용정보 **외부 미전송**(`AI_EXTERNAL_CALL_ENABLED=false`)",
 "LLM을 안 부르는 건 기술 부족이 아니라 **금융권 클라우드·신용정보 국외이전 기준을 코드로 준수한 설계 결정** — 파일럿에서 행내 LLM으로 교체할 인터페이스(`narrate.py`·`llm.py`) 보존",
], size=13, gap=9)
footer(s)

# ══════════════════ 15. 실배포 ══════════════════
s = slide()
headline(s, "실배포 — 연동 없이 시작하는 파일럿", "연동은 파일럿의 전제가 아니라 결과물이다")
table(s, MARG, 1.75, 12.2, [
 ["단계", "기간", "내용", "연동"],
 ["파일럿 전", "4주", "준법 문구 승인 · 상품 조건 검증 · 모델 대장 · PIA · KB 고시 수기 반입", "0"],
 ["P1 섀도", "8주", "RM 단독 — 화면 판정 vs RM 판단 비교 기록", "0"],
 ["P2 배석", "8주", "동의 고객 RM 배석 열람 → 이해도·설명의무 측정", "0"],
 ["연동 후", "예산 승인 시", "SSO · 원장/여신 · CRM · 실 OCR · 서버 SLA", "IT 예산"],
], [1.7, 1.8, 7.2, 1.5], size=13, row_h=0.5)
bullets(s, MARG, 4.9, 12.2, 1.7, [
 "측정 도구를 화면에 내장 — **RM 판정 기록(유지/변경+사유 → CSV)** = 추천 유지율 원자료",
 "**Go 기준** — 추천 유지율 ≥70% · 준법 오류 0 · 차단 판정 타당성 ≥90% · 고객 이해도 ≥70%",
 "못 재는 것은 과대약속 안 함 — 상담 전환율·준비시간 단축은 연동 후에만 (성공 기준에서 제외)",
], size=14, gap=9)
footer(s)

# ══════════════════ 16. 입상 7요건 ══════════════════
s = slide()
headline(s, "입상 7요건 충족표", "말만이 아니라 테스트로 지킨다")
table(s, MARG, 1.7, 12.2, [
 ["#", "요건", "충족"],
 ["1", "자체 정량 지표", "**BBP** — 개인화·팻테일·개입보정 (GARCH 아님)"],
 ["2", "설명가능성 XAI", "FX-EWI V/J/C/M 분해 + BBP 근거·출처"],
 ["3", "멀티에이전트", "Sentinel / Analyst / Advisor / Hedge + AI 모드(대화)"],
 ["4", "선제성", "BBP 임계 선제 경보 + 반복결제 계약 전 진단"],
 ["5", "KB 실배포", "KB Star FX 증강 + KB 공식상품 라우팅 + RM 티켓"],
 ["6", "검증", "walk-forward / PBO / DSR 정직 수행 + 한계 명시"],
 ["7", "차별화", "#6 블루오션 + BBP 개인화 + 정직 검증 + KB 통합"],
], [0.6, 3.0, 8.6], size=13, row_h=0.44)
rect(s, MARG, 5.6, 12.2, 0.95, BG)
txt(s, 0.85, 5.76, 11.6, 0.7,
    "검증 현황(재현 가능) — 데모 UI 스모크 **368** · 서버 pytest **257** · 문서 정합 **16** · 텍스트 확대(WCAG)·모바일 무회귀 — **전부 통과**",
    size=14.5, color=INK)
footer(s)

# ══════════════════ 17. 마무리 ══════════════════
s = slide()
rect(s, 0, 0, PAGE_W, PAGE_H, NAVY)
rect(s, 0, 3.62, PAGE_W, 0.045, GOLD)
txt(s, 1.0, 1.5, 11.3, 0.9, "예측 알파는 죽었지만,", size=28, color=RGBColor(0xC9, 0xD2, 0xE2))
txt(s, 1.0, 2.15, 11.3, 1.0, "더 정직하고 SME에 유용한 서비스로 재무장했습니다.", size=30, color=WHITE, bold=True)
txt(s, 1.0, 4.0, 11.3, 0.8,
    "회사별 예산환율 리스크 번역(BBP)  +  정직한 검증  +  KB 공식상품 원스톱 라우팅",
    size=19, color=GOLD, bold=True)
txt(s, 1.0, 5.0, 11.3, 0.6,
    "경진대회는 퀀트펀드가 아니라 KB 서비스 심사 — 이 방향이 입상 요건을 정직하게 지킵니다.",
    size=14.5, color=RGBColor(0x9A, 0xA7, 0xBD))
txt(s, 1.0, 5.9, 11.3, 0.4, "▶ https://seonu-dragon.github.io/KB_FX_Sentinel/", size=14, color=GOLD, bold=True)

out = os.path.join(ROOT, "FX_Sentinel_기술설명서.pptx")
prs.save(out)
print("saved:", out, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
